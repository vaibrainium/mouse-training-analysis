import datetime
import pickle
import tempfile
from pathlib import Path

import matplotlib
import numpy as np
import pandas as pd
from matplotlib import cm, gridspec
from matplotlib import pyplot as plt
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from src.utils import pmf_utils

# Constants
TEXTWIDTH = 5.0  # 9.0  # inches
FONTSIZE = 5.0  # 9.0

matplotlib.rcParams.update(
    {
        "font.size": FONTSIZE,
        "axes.titlesize": FONTSIZE,
        "axes.labelsize": FONTSIZE,
        "xtick.labelsize": FONTSIZE,
        "ytick.labelsize": FONTSIZE,
        "legend.fontsize": FONTSIZE,
        "figure.titlesize": FONTSIZE,
        "figure.dpi": 300,
        "savefig.dpi": 300,
    }
)

color = ["#d11149", "#1a8fe3", "#1ccd6a", "#e6c229", "#6610f2", "#f17105", "#65e5f3", "#bd8ad5", "#b16b57"]


def get_recent_sessions(data_dir=Path("Y:/"), last_X_business_days=None, start_date=None, end_date=None):
    # Ensure data_dir is a Path object
    data_dir = Path(data_dir)

    # Get all mouse ID folders except "XXX"
    mouse_ids = [f for f in data_dir.iterdir() if f.is_dir() and f.name != "XXX"]

    # Determine date range
    today = pd.to_datetime("today")  # Keep full timestamp
    if last_X_business_days is not None:
        start_date = today - pd.offsets.BDay(last_X_business_days)
        end_date = today
    else:
        start_date = pd.to_datetime(start_date, errors="coerce")
        end_date = pd.to_datetime(end_date, errors="coerce") if end_date else today

    # Validate date range
    if pd.isna(start_date) or pd.isna(end_date):
        raise ValueError("Invalid start_date or end_date format.")

    # Convert start_date and end_date to a consistent format
    date_format = "%Y-%m-%d"
    start_date_str = start_date.strftime(date_format)
    end_date_str = end_date.strftime(date_format)

    session_data = []
    required_columns = ["date", "start_weight", "end_weight", "baseline_weight", "experiment", "session"]

    for mouse in mouse_ids:
        history_path = mouse / "history.csv"

        try:
            history = pd.read_csv(history_path)
        except FileNotFoundError:
            continue

        # Convert date column dynamically to handle mixed formats
        history["date"] = pd.to_datetime(history["date"], format="mixed", errors="coerce")
        history = history.dropna(subset=["date"])  # Drop rows with invalid dates

        # Reformat history["date"] to match start_date and end_date format
        history["date"] = history["date"].dt.strftime(date_format)
        history["date"] = pd.to_datetime(history["date"], format=date_format, errors="coerce")  # Convert back to datetime

        # Apply date filtering
        history = history[(history["date"] >= start_date) & (history["date"] <= end_date)]

        if history.empty:
            continue

        # Ensure required columns exist
        history = history.reindex(columns=required_columns, fill_value=pd.NA)
        history["start_weight"] = history["start_weight"] / history["baseline_weight"] * 100
        history["end_weight"] = history["end_weight"] / history["baseline_weight"] * 100
        history["mouse_id"] = mouse.name

        session_data.append(history[["mouse_id", "date", "start_weight", "end_weight", "experiment", "session"]])

    # Concatenate all results into a single DataFrame
    return pd.concat(session_data, ignore_index=True) if session_data else pd.DataFrame()


def preprocess_data(data):
    data = data.dropna(subset=["outcome"])
    data = data[data.is_correction_trial == False]
    return data


def get_binned_accuracy(data, bin_size=10):
    outcome_array = data["outcome"].astype(float).to_numpy()
    num_bins = len(outcome_array) // bin_size
    binned_accuracy = [np.nanmean(outcome_array[i * bin_size : (i + 1) * bin_size]) * 100 for i in range(num_bins)]
    binned_indices = np.arange(num_bins) * bin_size
    return binned_indices, np.array(binned_accuracy)


def plot_binned_session_accuracy(data, bin_size, color, ax, label=None, tittle=None):
    binned_indices, binned_accuracy = get_binned_accuracy(data, bin_size=bin_size)
    ax.plot(binned_indices, binned_accuracy, color=color, linewidth=2, alpha=0.7, label=label)
    ax.set_ylim(0, 100)
    ax.set_xlabel("Trial Number")
    ax.set_ylabel("Accuracy (%)")
    ax.set_title(tittle, fontsize=FONTSIZE + 3)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)


def plot_psychometric_function(data, color, ax, label=None, tittle=None):
    x_data, y_data, _, x_model, y_model = pmf_utils.get_psychometric_data(data)
    ax.plot(x_data, y_data, "o", color=color)
    ax.plot(x_model, y_model, color=color, label=label)
    ax.set_ylim(0, 1)
    ax.set_xlabel("Coherence (%)")
    ax.set_ylabel("Prop. of Positive Choices (%)")
    ax.set_title(tittle, fontsize=FONTSIZE + 3)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)


def plot_accuracy_function(data, color, ax, label=None, tittle=None):
    x_data, y_data = pmf_utils.get_accuracy_data(data)
    ax.plot(x_data, y_data * 100, "o-", color=color, label=label)
    ax.set_ylim(0, 100)
    ax.set_xlabel("Coherence (%)")
    ax.set_ylabel(" % Correct")
    ax.set_title(tittle, fontsize=FONTSIZE + 3)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)


def plot_chronometric_function(data, color, ax, label=None, tittle=None):
    x_data, y_data = pmf_utils.get_chronometric_data(data)

    ax.plot(x_data, y_data, "o-", color=color, label=label)
    ax.set_ylim(0, 3)
    ax.set_xlabel("Coherence (%)")
    ax.set_ylabel("Reaction Time (ms)")
    ax.set_title(tittle, fontsize=FONTSIZE + 3)
    ax.spines["top"].set_visible(False)
    ax.spines["right"].set_visible(False)


def add_to_powerpoint(fig, mouse_id, prs):
    """
    Adds a matplotlib figure to a new slide in a PowerPoint presentation.

    Parameters:
    - fig: Matplotlib figure object.
    - mouse_id: Identifier for the mouse (used as the slide title).
    - pptx_filename: Path to the PowerPoint file where the slide will be added.
    """
    # Save figure to a temporary file
    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmpfile:
        fig.savefig(tmpfile.name, dpi=300, bbox_inches="tight")
        img_path = tmpfile.name  # Get the saved image path

    # Add a new slide
    slide_layout = prs.slide_layouts[5]  # Title only layout
    slide = prs.slides.add_slide(slide_layout)

    # Add title with mouse_id
    title = slide.shapes.title
    title.text = f"Mouse ID: {mouse_id}"
    title.text_frame.paragraphs[0].font.size = Pt(24)  # Adjust title font size
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Add image to slide
    left = Inches(0.8)
    top = Inches(1.2)
    slide.shapes.add_picture(img_path, left, top, width=Inches(9))

    print(f"Added {mouse_id}")


data_dir = Path("Y:/")
pptx_file1 = f"Y:/daily-data/{datetime.datetime.now().strftime('%Y-%m-%d')}.pptx"
pptx_file2 = f"D:/Data/research/mouse-training/daily-data/{datetime.datetime.now().strftime('%Y-%m-%d')}.pptx"

prs = Presentation()

session_info = get_recent_sessions(data_dir=data_dir, last_X_business_days=2)

session_info.date = pd.to_datetime(session_info.date).dt.date

# sort by mouse_id
session_info = session_info.sort_values(by="mouse_id", ascending=True)
for mouse_id in session_info.mouse_id.unique():
    mouse_sessions = session_info[session_info.mouse_id == mouse_id]
    # sort by date
    mouse_sessions = mouse_sessions.sort_values(by="date", ascending=True)

    fig, ax = plt.subplots(2, 3, figsize=(1.5 * TEXTWIDTH, TEXTWIDTH), dpi=300, constrained_layout=True)
    for idx_date, date in enumerate(mouse_sessions.date.unique()):
        sessions = mouse_sessions[mouse_sessions.date == date].reset_index()

        for idx, metadata in sessions.iterrows():
            trial_info = pd.read_csv(
                data_dir / mouse_id / "data/random_dot_motion" / metadata.experiment / metadata.session / f"{mouse_id}_trial.csv"
            )
            trial_info = preprocess_data(trial_info)

            plot_binned_session_accuracy(
                data=trial_info,
                bin_size=25,
                color=color[idx],
                ax=ax[idx_date, 0],
                label=f"Session -{idx+1}",
                tittle=f"Binned Accuracy \n ({metadata.date})",
            )
            plot_accuracy_function(
                data=trial_info, color=color[idx], ax=ax[idx_date, 1], label=f"Session {idx+1}", tittle=f"Accuracy \n ({metadata.date})"
            )
            plot_chronometric_function(
                data=trial_info, color=color[idx], ax=ax[idx_date, 2], label=f"Session {idx+1}", tittle=f"Median Reaction Time \n ({metadata.date})"
            )

            ax[idx_date, -1].legend(loc="upper left", bbox_to_anchor=(1, 1))
    # plt.show()
    add_to_powerpoint(fig, mouse_id, prs)

# Save PowerPoint file
prs.save(pptx_file1)
prs.save(pptx_file2)
